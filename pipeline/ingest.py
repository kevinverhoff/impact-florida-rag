"""
Step 3: Text Extractor

Reads metadata.json, extracts text from every downloaded file, and writes
documents.parquet -- one row per document.

Output columns (all metadata fields carried through, plus):
  text              Full extracted text
  headings          JSON string: list of {level, text, char_offset} (DOCX/XLSX/PPTX)
  char_count        len(text)
  extraction_status "ok" | "error" | "skipped"
  extraction_error  Error message, or null

Dispatches on file extension (not original MIME type) so that exported
Google-native files (.docx, .xlsx, .pptx) are handled by the right extractor.

Usage:
  python pipeline/ingest.py
"""

import io
import json
import logging
import tempfile
import traceback
from pathlib import Path
from typing import Any

import pandas as pd
import pdfplumber
import docx
from pptx import Presentation
from openpyxl import load_workbook

# pdfminer emits font-descriptor warnings on some PDFs that do not affect
# text extraction quality. Suppress them so they do not flood the console.
logging.getLogger("pdfminer").setLevel(logging.ERROR)

PROJECT_ROOT = Path(__file__).parent.parent
METADATA_PATH = PROJECT_ROOT / "data" / "metadata.json"
OUTPUT_PATH = PROJECT_ROOT / "data" / "documents.parquet"

# File extensions that cannot yield useful text
SKIP_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".mp4", ".mp3", ".zip", ".bin"}

# ---------------------------------------------------------------------------
# Extractors
# ---------------------------------------------------------------------------

def _clean(text: str) -> str:
    """Collapse excessive blank lines and strip leading/trailing whitespace."""
    lines = text.splitlines()
    cleaned, prev_blank = [], False
    for line in lines:
        is_blank = not line.strip()
        if is_blank and prev_blank:
            continue
        cleaned.append(line)
        prev_blank = is_blank
    return "\n".join(cleaned).strip()


def extract_pdf(path: Path) -> tuple[str, list[dict], str, str | None]:
    """Extract text page-by-page. No heading detection for PDFs."""
    pages = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                pages.append(page_text)
    text = _clean("\n\n".join(pages))
    return text, [], "ok", None


def extract_docx(path: Path) -> tuple[str, list[dict], str, str | None]:
    """
    Extract paragraphs and tables from a Word document.
    Detects Heading 1/2/3 styles and records their character offsets so
    the chunker can split at section boundaries.
    """
    doc = docx.Document(path)
    parts: list[str] = []
    headings: list[dict] = []
    current_offset = 0

    for para in doc.paragraphs:
        para_text = para.text.strip()
        if not para_text:
            continue

        style_name = para.style.name if para.style else ""
        for level in [1, 2, 3]:
            if f"Heading {level}" in style_name:
                headings.append({
                    "level": level,
                    "text": para_text,
                    "char_offset": current_offset,
                })
                break

        parts.append(para_text)
        current_offset += len(para_text) + 1  # +1 for the joining newline

    # Tables
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)
                current_offset += len(row_text) + 1

    text = _clean("\n".join(parts))
    return text, headings, "ok", None


def extract_xlsx(path: Path) -> tuple[str, list[dict], str, str | None]:
    """
    Extract cell values sheet-by-sheet. Sheet names become pseudo-headings
    so the chunker can split at sheet boundaries.
    """
    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    headings: list[dict] = []
    current_offset = 0

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_header = f"[Sheet: {sheet_name}]"
        headings.append({
            "level": 1,
            "text": sheet_header,
            "char_offset": current_offset,
        })
        parts.append(sheet_header)
        current_offset += len(sheet_header) + 1

        for row in ws.iter_rows(values_only=True):
            row_values = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if not row_values:
                continue
            row_text = "\t".join(row_values)
            parts.append(row_text)
            current_offset += len(row_text) + 1

    wb.close()
    text = _clean("\n".join(parts))
    return text, headings, "ok", None


def extract_pptx(path: Path) -> tuple[str, list[dict], str, str | None]:
    """
    Extract text slide-by-slide. Title shapes become headings.
    Per-shape errors are skipped so a single bad shape does not abort the file.
    """
    prs = Presentation(path)
    parts: list[str] = []
    headings: list[dict] = []
    current_offset = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_title = None

        for shape in slide.shapes:
            try:
                if not shape.has_text_frame:
                    continue
                shape_text = shape.text_frame.text.strip()
                if not shape_text:
                    continue

                # Placeholder index 0 is the title
                is_title = (
                    hasattr(shape, "placeholder_format")
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.idx == 0
                )
                if is_title:
                    slide_title = shape_text
                    headings.append({
                        "level": 1,
                        "text": shape_text,
                        "char_offset": current_offset,
                    })

                parts.append(shape_text)
                current_offset += len(shape_text) + 1
            except Exception:
                # Skip malformed shapes -- do not abort the whole file
                continue

        if slide_title is None:
            # Anchor every slide even if it has no title placeholder
            fallback = f"[Slide {slide_num}]"
            headings.append({"level": 1, "text": fallback, "char_offset": current_offset})

    text = _clean("\n".join(parts))
    return text, headings, "ok", None


def extract_csv(path: Path) -> tuple[str, list[dict], str, str | None]:
    """Read CSV and render as tab-separated text with a column-header heading."""
    df = pd.read_csv(path, dtype=str).fillna("")
    header = "\t".join(df.columns.tolist())
    headings = [{"level": 1, "text": header, "char_offset": 0}]
    rows = [header] + ["\t".join(row) for row in df.values.tolist()]
    text = _clean("\n".join(rows))
    return text, headings, "ok", None


def extract_txt(path: Path) -> tuple[str, list[dict], str, str | None]:
    text = _clean(path.read_text(encoding="utf-8", errors="replace"))
    return text, [], "ok", None


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

_EXTRACTORS = {
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".doc":  extract_docx,   # python-docx handles many .doc files
    ".xlsx": extract_xlsx,
    ".xls":  extract_xlsx,
    ".pptx": extract_pptx,
    ".csv":  extract_csv,
    ".txt":  extract_txt,
}


def extract(local_path: Path, original_mime: str) -> tuple[str, list[dict], str, str | None]:
    """
    Dispatch to the correct extractor based on file extension.
    Returns (text, headings, status, error_message).
    """
    ext = local_path.suffix.lower()

    if ext in SKIP_EXTENSIONS:
        return "", [], "skipped", None

    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        return "", [], "skipped", f"No extractor for extension: {ext}"

    if not local_path.exists():
        return "", [], "error", "File not found on disk"

    try:
        return extractor(local_path)
    except Exception:
        return "", [], "error", traceback.format_exc(limit=3)

# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    if not METADATA_PATH.exists():
        raise FileNotFoundError("metadata.json not found -- run get_docs.py first")

    with open(METADATA_PATH, encoding="utf-8") as f:
        metadata: list[dict] = json.load(f)

    processable = [
        m for m in metadata
        if m.get("download_status") in ("downloaded", "exists")
    ]
    skipped_download = len(metadata) - len(processable)
    print(f"Extracting text from {len(processable)} files ({skipped_download} skipped/errored in download)...")

    rows: list[dict[str, Any]] = []
    counts: dict[str, int] = {"ok": 0, "error": 0, "skipped": 0}

    for idx, record in enumerate(processable, 1):
        local_path = PROJECT_ROOT / record["local_path"]
        original_mime = record.get("mime_type", "")

        text, headings, status, error = extract(local_path, original_mime)

        if status == "error":
            print(f"  [{idx}/{len(processable)}] ERROR  {record['file_name'][:60]}")
            first_line = error.splitlines()[0] if error else "unknown"
            print(f"    {first_line}")

        counts[status] = counts.get(status, 0) + 1

        rows.append({
            **{k: record.get(k) for k in (
                "file_id", "file_name", "mime_type", "folder_path",
                "local_path", "drive_url",
                "program", "doc_type", "academic_year", "season",
                "date_precision", "district",
            )},
            "text": text,
            "headings": json.dumps(headings, ensure_ascii=False),
            "char_count": len(text),
            "extraction_status": status,
            "extraction_error": error,
        })

    df = pd.DataFrame(rows)
    df.to_parquet(OUTPUT_PATH, index=False)

    ok_rows = df[df["extraction_status"] == "ok"]
    print(f"\nDone.")
    print(f"  OK:       {counts.get('ok', 0)}")
    print(f"  Skipped:  {counts.get('skipped', 0)}")
    print(f"  Errors:   {counts.get('error', 0)}")
    if len(ok_rows):
        print(f"  Avg text: {int(ok_rows['char_count'].mean()):,} chars")
    print(f"  Output:   {OUTPUT_PATH}")




# ---------------------------------------------------------------------------
# Streaming extraction (no disk read -- accepts BytesIO)
# ---------------------------------------------------------------------------

def extract_from_bytes(
    content: "io.BytesIO | None", ext: str
) -> tuple[str, list[dict], str, "str | None"]:
    """
    Extract text from an in-memory BytesIO object.

    Writes content to a temporary file, delegates to the matching extractor,
    then deletes the temp file immediately. No data is persisted.
    """
    if content is None:
        return "", [], "skipped", "No content (download failed)"
    if ext in SKIP_EXTENSIONS:
        return "", [], "skipped", None
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        return "", [], "skipped", f"No extractor for extension: {ext}"

    content.seek(0)
    suffix = ext if ext.startswith(".") else f".{ext}"
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
            tmp.write(content.read())
            tmp_path = Path(tmp.name)
        return extractor(tmp_path)
    except Exception:
        return "", [], "error", traceback.format_exc(limit=3)
    finally:
        if tmp_path and tmp_path.exists():
            tmp_path.unlink()


def process_stream(
    records_with_content,
    output_path: Path = OUTPUT_PATH,
) -> None:
    """
    Consume a (metadata_dict, BytesIO, ext) generator and write documents.parquet.

    Memory-efficient: each BytesIO is extracted and discarded before the next
    file is downloaded, so the full corpus is never in RAM simultaneously.
    """
    rows: list[dict] = []
    counts: dict[str, int] = {"ok": 0, "error": 0, "skipped": 0}

    for idx, (record, buf, ext) in enumerate(records_with_content, 1):
        text, headings, status, error = extract_from_bytes(buf, ext)

        if status == "error":
            fname = record.get("file_name", "")[:60]
            print(f"  [{idx}] ERROR  {fname}")
            if error:
                print(f"    {error.splitlines()[0]}")

        counts[status] = counts.get(status, 0) + 1
        rows.append({
            **{k: record.get(k) for k in (
                "file_id", "file_name", "mime_type", "folder_path",
                "local_path", "drive_url",
                "program", "doc_type", "academic_year", "season",
                "date_precision", "district",
            )},
            "text":              text,
            "headings":          json.dumps(headings, ensure_ascii=False),
            "char_count":        len(text),
            "extraction_status": status,
            "extraction_error":  error,
        })

    df = pd.DataFrame(rows)
    df.to_parquet(output_path, index=False)

    ok_rows = df[df["extraction_status"] == "ok"]
    print(f"\nDone (streaming).")
    print(f"  OK:       {counts.get('ok', 0)}")
    print(f"  Skipped:  {counts.get('skipped', 0)}")
    print(f"  Errors:   {counts.get('error', 0)}")
    if len(ok_rows):
        print(f"  Avg text: {int(ok_rows['char_count'].mean()):,} chars")
    print(f"  Output:   {output_path}")
if __name__ == "__main__":
    main()